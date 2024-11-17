import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from datetime import datetime, timedelta
import os
import MetaTrader5 as mt5
from mplfinance.original_flavor import candlestick_ohlc
from mpl_toolkits.axes_grid1 import make_axes_locatable
import re
from pptx import Presentation
from pptx.util import Inches

# Initialize MT5 connection
if not mt5.initialize():
    print("MT5 initialization failed")
    mt5.shutdown()

def estimate_typical_spread(symbol):
    major_fx = {
        "EURUSD.sd": 0.00010, "USDJPY.sd": 0.00010, "GBPUSD.sd": 0.00015,
        "USDCHF.sd": 0.00015, "AUDUSD.sd": 0.00015, "USDCAD.sd": 0.00015
    }
    minor_fx = {
        "EURGBP.sd": 0.00020, "EURJPY.sd": 0.00020, "GBPJPY.sd": 0.00025,
        "EURCHF.sd": 0.00025, "AUDJPY.sd": 0.00020, "NZDUSD.sd": 0.00020
    }
    
    if symbol in major_fx:
        return major_fx[symbol]
    elif symbol in minor_fx:
        return minor_fx[symbol]
    else:
        return 0.00030  # Default to 3 pips

def get_current_spread(symbol):
    symbol_info = mt5.symbol_info(symbol)
    if symbol_info is None:
        print(f"Failed to get symbol info for {symbol}")
        return None
    
    spread = symbol_info.ask - symbol_info.bid
    if spread == 0:
        typical_spread = estimate_typical_spread(symbol)
        print(f"MT5 returned 0 spread for {symbol}. Using estimated typical spread: {typical_spread}")
        return typical_spread
    
    return spread

def fetch_data(symbol, start_time, end_time):
    rates = mt5.copy_rates_range(symbol, mt5.TIMEFRAME_M1, start_time, end_time)
    if rates is None or len(rates) == 0:
        print(f"No OHLC data available for {symbol}")
        return None

    df = pd.DataFrame(rates)
    df['time'] = pd.to_datetime(df['time'], unit='s')
    df = calculate_atr(df)

    # Check if the 'volume' column exists, and if not, add a placeholder
    if 'volume' not in df.columns:
        df['volume'] = 0

    if df.empty:
        print(f"DataFrame is empty after filtering for {symbol}")
        return None

    return df

def calculate_atr(df, period=14):
    df['tr1'] = df['high'] - df['low']
    df['tr2'] = (df['high'] - df['close'].shift()).abs()
    df['tr3'] = (df['low'] - df['close'].shift()).abs()
    df['tr'] = df[['tr1', 'tr2', 'tr3']].max(axis=1)
    df['atr'] = df['tr'].rolling(window=period).mean()
    return df

def parse_excel_input(excel_input):
    lines = excel_input.strip().split('\n')
    events = []
    date_time_inputs = {}
    event_data_list = {}
    event_symbols = {}

    current_event_name = ""
    current_event_time = ""

    for line in lines:
        if line.startswith('Event'):
            parts = line.split('\t')
            current_event_name = parts[1].strip()
            events.append({'name': current_event_name})
        elif line.startswith('Actual:'):
            parts = line.split('\t')
            symbol = parts[-1].strip()
            event_symbols[current_event_name] = {"mt5": symbol}
        elif line.startswith('Forecast:'):
            # Ignore, as we're using MT5 data only
            pass
        elif line.startswith('Time (GMT):'):
            parts = line.split('\t')
            current_event_time = parts[1].strip()
        elif re.match(r'\d{2}-\w{3}-\d{2}', line):
            parts = line.split('\t')
            date_str = parts[0].strip()
            date = datetime.strptime(date_str, '%d-%b-%y').date()
            actual = parts[1].strip().replace('%', '')
            forecast = parts[2].strip().replace('%', '') or "NA"
            previous = parts[3].strip().replace('%', '') or "NA"
            time_str = parts[4].strip()

            if date not in date_time_inputs:
                date_time_inputs[date] = {
                    "date": date,
                    "event": current_event_time + ":00",
                }
                event_data_list[date] = []

            event_data_list[date].append({
                "name": current_event_name,
                "actual": actual,
                "forecast": forecast,
                "previous": previous,
                "time": time_str
            })

    return date_time_inputs, event_data_list, events, event_symbols

# Initialize presentation
template_path = "/Users/pawan/Desktop/Quantwater Tech Investments/Research & Development/Research/Blog/Slides/template.pptx"
prs = Presentation(template_path)
generated_image_filenames = []

def plot_ohlc(df, symbol, event_time, event_data, atr_multiple, spread, atr_period=14,
              is_last_chart=False, event_summary=None, start_time=None, end_time=None):
    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(14.5, 9),
                                   gridspec_kw={'height_ratios': [4, 1]}, sharex=True)
    
    df['time_num'] = mdates.date2num(df['time'])
    ohlc = df[['time_num', 'open', 'high', 'low', 'close']].values

    plt.style.use('ggplot')

    candlestick_ohlc(ax1, ohlc, width=0.0004, colorup='green', colordown='red')

    ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M'))
    plt.setp(ax1.get_xticklabels(), rotation=45, ha='right')

    closest_index = (df['time'] - event_time).abs().idxmin()
    closest_time = df.loc[closest_index, 'time']
    event_price = df.loc[closest_index, 'open']

    ax1.plot(closest_time, event_price, 'X', color='red', markersize=10, label='Event Time')

    atr_value = df.loc[closest_index, 'atr']
    volatility_unit = atr_value * atr_multiple + spread

    for i in range(1, 6):
        upper_range = event_price + i * volatility_unit
        lower_range = event_price - i * volatility_unit
        ax1.axhline(upper_range, color='blue', linestyle='--', label=f'+{i} Volatility Unit' if i == 1 else "")
        ax1.axhline(lower_range, color='orange', linestyle='--', label=f'-{i} Volatility Unit' if i == 1 else "")

    post_event_df = df[df['time'] >= event_time]
    max_price = post_event_df['high'].max()
    min_price = post_event_df['low'].min()
    max_movement = max_price - event_price
    min_movement = event_price - min_price

    if not post_event_df[post_event_df['high'] == max_price].empty:
        max_time = post_event_df.loc[post_event_df['high'] == max_price, 'time'].iloc[0]
    else:
        max_time = None  # No match found for max price

    if not post_event_df[post_event_df['low'] == min_price].empty:
        min_time = post_event_df.loc[post_event_df['low'] == min_price, 'time'].iloc[0]
    else:
        min_time = None  # No match found for min price

    if max_time:
        ax1.plot(max_time, max_price, 'gx', markersize=10, label='Max Peak')
    if min_time:
        ax1.plot(min_time, min_price, 'rx', markersize=10, label='Min Peak')

    pre_atr_df = df[df['time'] < event_time].iloc[-(atr_period*2):-atr_period]
    pre_atr_volatility = pre_atr_df['high'].max() - pre_atr_df['low'].min()

    volume_color = 'blue'
    volume_alpha = 0.3
    ax2.bar(df['time'], df['volume'], width=0.0004, color=volume_color, alpha=volume_alpha)
    ax2.set_ylabel('Volume')

    info_text = (
        f'ATR: {atr_value:.4f}\n'
        f'Spread: {spread:.5f}\n'
        f'Pre-ATR Volatility: {pre_atr_volatility:.4f}\n'
        f'\n'
        f'Volatility Unit (ATR*{atr_multiple}+spread): {volatility_unit:.4f}\n'
        f'Max Movement: {max_movement/volatility_unit:.2f}x\n'
        f'Min Movement: {min_movement/volatility_unit:.2f}x'
    )

    if is_last_chart and event_summary and event_summary['Avg Volatility Unit'] > 0:
        info_text += (
            f'\n\nEvent Summary:\n'
            f'Avg Volatility Unit: {event_summary["Avg Volatility Unit"]:.4f}\n'
            f'Avg Max Movement: {event_summary["Avg Max Movement"]:.2f}x\n'
            f'Avg Min Movement: {event_summary["Avg Min Movement"]:.2f}x\n'
            f'Avg Pre-ATR Volatility: {event_summary["Avg Pre-ATR Volatility"]:.4f}'
        )

    ax1.text(1.01, 0.5, info_text, transform=ax1.transAxes, ha='left', va='center', fontsize=9)

    if start_time and end_time:
        ax1.set_xlim(mdates.date2num(start_time), mdates.date2num(end_time))

    divider = make_axes_locatable(ax2)
    ax_event = divider.append_axes("bottom", size="15%", pad=0.1, sharex=ax1)
    ax_event.axis('off')

    plt.subplots_adjust(bottom=0.45)  

    for i, event in enumerate(event_data):
        event_name = event['name']
        actual = event['actual']
        forecast = event['forecast']
        previous = event['previous']

        if forecast == "NA":
            event_color = 'black'
        else:
            try:
                if float(actual) > float(forecast):
                    event_color = 'green'
                elif float(actual) < float(forecast):
                    event_color = 'red'
                else:
                    event_color = 'blue'
            except ValueError:
                event_color = 'black'

        event_time_num = mdates.date2num(event_time)
        event_text_line = f"{event_name} | Actual: {actual} | Forecast: {forecast} | Previous: {previous}"
        y_offset = -0.2 - i * 0.1  
        ax_event.text(event_time_num, y_offset, event_text_line, color=event_color,
                      ha='center', va='top', fontsize=9, rotation=0, transform=ax_event.transData, fontweight ='bold')

    ax1.set_title(f'{symbol} Price Action - {event_time.strftime("%Y-%m-%d %H:%M")} GMT', fontweight ='bold')
    ax1.set_ylabel('Price')
    ax1.grid(True)
    ax1.legend(loc='upper left', bbox_to_anchor=(1, 1))

    plt.tight_layout()
    plt.subplots_adjust(hspace=0)

    output_dir = "output_plots"
    os.makedirs(output_dir, exist_ok=True)
    image_filename = os.path.join(output_dir, f"{symbol}_{event_time.strftime('%Y%m%d_%H%M%S')}.png")


    plt.savefig(image_filename)
    plt.close(fig)

    return image_filename, volatility_unit, max_movement, min_movement, pre_atr_volatility

def create_pptx(generated_image_filenames):
    output_dir = "output_plots"
    prs = Presentation()

    for image_filename in generated_image_filenames:
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Add a new slide
        full_image_path = os.path.join(output_dir, image_filename)
        pic = slide.shapes.add_picture(full_image_path, Inches(0), Inches(0), width=Inches(16), height=Inches(9))

    output_pptx_path = "output.pptx"
    prs.save(output_pptx_path)

    return output_pptx_path, generated_image_filenames

def process_event(event, date_time_inputs, event_data_list, event_symbols):
    event_summary = {
        'volatility_units': [],
        'max_movements': [],
        'min_movements': [],
        'pre_atr_volatilities': [],
        'Avg Volatility Unit': 0,
        'Avg Max Movement': 0,
        'Avg Min Movement': 0,
        'Avg Pre-ATR Volatility': 0
    }

    event_dates = [
        date for date in sorted(date_time_inputs.keys())
        if any(e['name'] == event['name'] for e in event_data_list[date])
    ]

    for date in event_dates:
        event_time = datetime.strptime(date_time_inputs[date]["event"], '%H:%M:%S').time()
        event_datetime = datetime.combine(date_time_inputs[date]["date"], event_time)
        event_data = event_data_list[date]

        for event_info in event_data:
            if event_info['name'] == event['name']:
                symbol = event_symbols[event['name']]["mt5"]
                df = fetch_data(symbol, event_datetime - timedelta(hours=2), event_datetime + timedelta(hours=2))

                if df is not None:
                    image_filename, volatility_unit, max_movement, min_movement, pre_atr_volatility = plot_ohlc(
                        df, symbol, event_datetime, event_data, atr_multiple=3, spread=get_current_spread(symbol),
                        is_last_chart=(event_info == event_data[-1]), event_summary=event_summary,
                        start_time=event_datetime - timedelta(hours=2), end_time=event_datetime + timedelta(hours=2)
                    )

                    event_summary['volatility_units'].append(volatility_unit)
                    event_summary['max_movements'].append(max_movement)
                    event_summary['min_movements'].append(min_movement)
                    event_summary['pre_atr_volatilities'].append(pre_atr_volatility)

                    print(f"Processed {symbol} at {event_datetime.strftime('%Y-%m-%d %H:%M')} GMT")
                    generated_image_filenames.append(image_filename)

    if event_summary['volatility_units']:
        event_summary['Avg Volatility Unit'] = sum(event_summary['volatility_units']) / len(event_summary['volatility_units'])
        event_summary['Avg Max Movement'] = sum(event_summary['max_movements']) / len(event_summary['max_movements'])
        event_summary['Avg Min Movement'] = sum(event_summary['min_movements']) / len(event_summary['min_movements'])
        event_summary['Avg Pre-ATR Volatility'] = sum(event_summary['pre_atr_volatilities']) / len(event_summary['pre_atr_volatilities'])
    else:
        print(f"No data available for event: {event['name']}")

    return event_summary

import sys

def main():
    atr_multiple = 1.5
    atr_period = 14

    if len(sys.argv) < 2:
        print("Usage: python analyzer-mt5-3.py '<excel_formatted_input>'")
        return

    excel_input = sys.argv[1]

    date_time_inputs, event_data_list, events, event_symbols = parse_excel_input(excel_input)

    generated_image_filenames = []

    for event in events:
        event_summary = process_event(event, date_time_inputs, event_data_list, event_symbols, generated_image_filenames)
        print(f"Event Summary for {event['name']}:")
        print(event_summary)

    output_pptx_path, _ = create_pptx(generated_image_filenames)

    # Clean up temporary image files
    for image_filename in generated_image_filenames:
        full_image_path = os.path.join("output_plots", image_filename)
        os.remove(full_image_path)

    # Shutdown MT5
    mt5.shutdown()

    print(f"Output PPTX file saved at: {output_pptx_path}")
    return output_pptx_path