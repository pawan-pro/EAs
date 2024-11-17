import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from datetime import datetime, timedelta
import os
import time
from polygon import RESTClient
import pytz
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import MetaTrader5 as mt5
from mplfinance.original_flavor import candlestick_ohlc
from mpl_toolkits.axes_grid1 import make_axes_locatable
import re
from pptx import Presentation
from pptx.util import Inches

# Polygon API key
API_KEY = "f8Oowzm4TdXsQjxOOlH7C1_2xB1dxNSB"

# Initialize Polygon RESTClient
client = RESTClient(api_key=API_KEY)

# Retry strategy to handle 429 errors
retry_strategy = Retry(
    total=2,
    backoff_factor=2,
    status_forcelist=[429, 500, 502, 503, 504],
    allowed_methods=["GET"]
)
adapter = HTTPAdapter(max_retries=retry_strategy)
http = requests.Session()
http.mount("https://", adapter)

# Initialize MT5 connection
if not mt5.initialize():
    print("MT5 initialization failed")
    mt5.shutdown()

def estimate_typical_spread(symbol):
    major_fx = {
        "EURUSD.sd": 0.00010, "USDJPY.sd": 0.00010, "GBPUSD.sd": 0.00015,
        "USDCHF.sd": 0.00015, "AUDUSD.sd": 0.00015, "USDCAD.sd": 0.00015
    }
    minor_fx = {
        "EURGBP.sd": 0.00020, "EURJPY.sd": 0.00020, "GBPJPY.sd": 0.00025,
        "EURCHF.sd": 0.00025, "AUDJPY.sd": 0.00020, "NZDUSD.sd": 0.00020
    }
    
    if symbol in major_fx:
        return major_fx[symbol]
    elif symbol in minor_fx:
        return minor_fx[symbol]
    else:
        return 0.00030  # Default to 3 pips

def get_current_spread(symbol):
    symbol_info = mt5.symbol_info(symbol)
    if symbol_info is None:
        print(f"Failed to get symbol info for {symbol}")
        return None
    
    spread = symbol_info.ask - symbol_info.bid
    if spread == 0:
        typical_spread = estimate_typical_spread(symbol)
        print(f"MT5 returned 0 spread for {symbol}. Using estimated typical spread: {typical_spread}")
        return typical_spread
    
    return spread

def fetch_data(symbol, start_time, end_time):
    if symbol == "USOILRoll":
        # Convert to UTC for MT5
        start_time_utc = start_time.astimezone(pytz.UTC) 

        end_time_utc = end_time.astimezone(pytz.UTC)
        
        rates = mt5.copy_rates_range(symbol, mt5.TIMEFRAME_M1, start_time_utc, end_time_utc)
        df = pd.DataFrame(rates)
        df['time'] = pd.to_datetime(df['time'], unit='s', utc=True)
        
        # Add volume column if it doesn't exist
        if 'volume' not in df.columns:
            df['volume'] = 0
            
        df = calculate_atr(df)
        return df
    else:
        # Fetch data from Polygon for other symbols
        api_symbol = symbol  
        aggs = []
        retry_attempts = 0
        current_date = start_time.astimezone(pytz.UTC).date()
        end_date = end_time.astimezone(pytz.UTC).date()

        while current_date <= end_date and retry_attempts < 5:
            try:
                for a in client.list_aggs(ticker=api_symbol, multiplier=1, timespan="minute", 
                                          from_=current_date.strftime('%Y-%m-%d'), 
                                          to=(current_date + timedelta(days=1)).strftime('%Y-%m-%d'), 
                                          limit=50000):
                    aggs.append(a)

                current_date += timedelta(days=1)
                retry_attempts = 0  
            except Exception as e:
                print(f"Error fetching data for {api_symbol} on {current_date}: {e}")
                if '429' in str(e):
                    retry_attempts += 1
                    sleep_time = 2 ** retry_attempts
                    print(f"Retrying in {sleep_time} seconds...")
                    time.sleep(sleep_time)
                else:
                    break

        if not aggs:
            print(f"No OHLC data available for {symbol}")
            return None

        df = pd.DataFrame([{
            'time': pd.to_datetime(agg.timestamp, unit='ms', utc=True),
            'open': agg.open,
            'high': agg.high,
            'low': agg.low,
            'close': agg.close,
            'volume': getattr(agg, 'volume', 0)  # Use 0 if volume doesn't exist
        } for agg in aggs])

        df = df[(df['time'] >= start_time.astimezone(pytz.UTC)) & 
                (df['time'] <= end_time.astimezone(pytz.UTC))]
        df.sort_values('time', inplace=True)
        df = calculate_atr(df)

        return df

def calculate_atr(df, period=14):
    df['tr1'] = df['high'] - df['low']
    df['tr2'] = (df['high'] - df['close'].shift()).abs()
    df['tr3'] = (df['low'] - df['close'].shift()).abs()
    df['tr'] = df[['tr1', 'tr2', 'tr3']].max(axis=1)
    df['atr'] = df['tr'].rolling(window=period).mean()
    return df

def parse_excel_input(excel_input):
    excel_input = excel_input.replace('^I', '\t')
    lines = excel_input.strip().split('\n')
    events = []
    date_time_inputs = {}
    event_data_list = {}
    event_symbols = {}

    current_event_name = ""
    current_event_time = ""

    for line in lines:
        if line.startswith('Event'):
            parts = line.split('\t')
            current_event_name = parts[1].strip()
            events.append({'name': current_event_name})
        elif line.startswith('Actual:'):
            parts = line.split('\t')
            symbol = parts[-1].strip()
            event_symbols[current_event_name] = {"polygon": symbol}
        elif line.startswith('Forecast:'):
            parts = line.split('\t')
            mt5_symbol = parts[-1].strip()
            event_symbols[current_event_name]["mt5"] = mt5_symbol
        elif line.startswith('Time (GMT):'):
            parts = line.split('\t')
            current_event_time = parts[1].strip()
        elif re.match(r'\d{2}-\w{3}-\d{2}', line):
            parts = line.split('\t')
            date_str = parts[0].strip()
            date = datetime.strptime(date_str, '%d-%b-%y').strftime('%Y-%m-%d')
            actual = parts[1].strip().replace('%', '')
            forecast = parts[2].strip().replace('%', '') or "NA"
            previous = parts[3].strip().replace('%', '') or "NA"
            time_str = parts[4].strip()

            if date not in date_time_inputs:
                date_time_inputs[date] = {
                    "date": date,
                    "event": current_event_time + "3:00",
                }
                event_data_list[date] = []

            event_data_list[date].append({
                "name": current_event_name,
                "actual": actual,
                "forecast": forecast,
                "previous": previous,
                "time": time_str
            })

    return date_time_inputs, event_data_list, events, event_symbols

# Initialize presentation
template_path = "/Users/pawan/Desktop/Quantwater Tech Investments/Research & Development/Research/Blog/Slides/template3.pptx"
prs = Presentation(template_path)
generated_image_filenames = []

def plot_ohlc(df, symbol, event_time, event_data, atr_multiple, spread, atr_period=14,
              is_last_chart=False, event_summary=None, start_time=None, end_time=None):
    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(14.5, 9),
                                   gridspec_kw={'height_ratios': [4, 1]}, sharex=True)
    
    df['time_num'] = mdates.date2num(df['time'])
    ohlc = df[['time_num', 'open', 'high', 'low', 'close']].values

    plt.style.use('ggplot')

    candlestick_ohlc(ax1, ohlc, width=0.0004, colorup='green', colordown='red')

    ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M'))
    plt.setp(ax1.get_xticklabels(), rotation=45, ha='right')

    closest_index = (df['time'] - event_time).abs().idxmin()
    closest_time = df.loc[closest_index, 'time']
    event_price = df.loc[closest_index, 'open']

    ax1.plot(closest_time, event_price, 'X', color='red', markersize=10, label='Event Time')

    atr_value = df.loc[closest_index, 'atr']
    volatility_unit = atr_value * atr_multiple + spread

    for i in range(1, 6):
        upper_range = event_price + i * volatility_unit
        lower_range = event_price - i * volatility_unit
        ax1.axhline(upper_range, color='blue', linestyle='--', label=f'+{i} Volatility Unit' if i == 1 else "")
        ax1.axhline(lower_range, color='orange', linestyle='--', label=f'-{i} Volatility Unit' if i == 1 else "")

    post_event_df = df[df['time'] >= event_time]
    max_price = post_event_df['high'].max()
    min_price = post_event_df['low'].min()
    max_movement = max_price - event_price
    min_movement = event_price - min_price

    if not post_event_df[post_event_df['high'] == max_price].empty:
        max_time = post_event_df.loc[post_event_df['high'] == max_price, 'time'].iloc[0]
        ax1.plot(max_time, max_price, 'gx', markersize=10, label='Max Peak')
    
    if not post_event_df[post_event_df['low'] == min_price].empty:
        min_time = post_event_df.loc[post_event_df['low'] == min_price, 'time'].iloc[0]
        ax1.plot(min_time, min_price, 'rx', markersize=10, label='Min Peak')

    pre_atr_df = df[df['time'] < event_time].iloc[-(atr_period*2):-atr_period]
    pre_atr_volatility = pre_atr_df['high'].max() - pre_atr_df['low'].min()

    # Only plot volume if it's not all zeros
    if df['volume'].any():
        volume_color = 'blue'
        volume_alpha = 0.3
        ax2.bar(df['time'], df['volume'], width=0.0004, color=volume_color, alpha=volume_alpha)
        ax2.set_ylabel('Volume')
    else:
        ax2.set_visible(False)
        plt.subplots_adjust(bottom=0.2)  # Adjust bottom margin when volume is hidden

    info_text = (
        f'ATR: {atr_value:.4f}\n'
        f'Spread: {spread:.5f}\n'
        f'Pre-ATR Volatility: {pre_atr_volatility:.4f}\n'
        f'\n'
        f'Volatility Unit (ATR*{atr_multiple}+spread): {volatility_unit:.4f}\n'
        f'Max Movement: {max_movement/volatility_unit:.2f}x\n'
        f'Min Movement: {min_movement/volatility_unit:.2f}x'
    )

    if is_last_chart and event_summary and event_summary['Avg Volatility Unit'] > 0:
        info_text += (
            f'\n\nEvent Summary:\n'
            f'Avg Volatility Unit: {event_summary["Avg Volatility Unit"]:.4f}\n'
            f'Avg Max Movement: {event_summary["Avg Max Movement"]:.2f}x\n'
            f'Avg Min Movement: {event_summary["Avg Min Movement"]:.2f}x\n'
            f'Avg Pre-ATR Volatility: {event_summary["Avg Pre-ATR Volatility"]:.4f}'
        )

    ax1.text(1.01, 0.5, info_text, transform=ax1.transAxes, ha='left', va='center', fontsize=9)

    if start_time and end_time:
        ax1.set_xlim(mdates.date2num(start_time), mdates.date2num(end_time))

    divider = make_axes_locatable(ax2)
    ax_event = divider.append_axes("bottom", size="15%", pad=0.1, sharex=ax1)
    ax_event.axis('off')

    plt.subplots_adjust(bottom=0.45)  

    for i, event in enumerate(event_data):
        event_name = event['name']
        actual = event['actual']
        forecast = event['forecast']
        previous = event['previous']

        if forecast == "NA":
            event_color = 'black'
        else:
            try:
                if float(actual) > float(forecast):
                    event_color = 'green'
                elif float(actual) < float(forecast):
                    event_color = 'red'
                else:
                    event_color = 'blue'
            except ValueError:
                event_color = 'black'

        event_time_num = mdates.date2num(event_time)
        event_text_line = f"{event_name} | Actual: {actual} | Forecast: {forecast} | Previous: {previous}"
        y_offset = -0.2 - i * 0.1  
        ax_event.text(event_time_num, y_offset, event_text_line, color=event_color,
                      ha='center', va='top', fontsize=9, rotation=0, transform=ax_event.transData, fontweight='bold')

    ax1.set_title(f'{symbol} Price Action - {event_time.strftime("%Y-%m-%d %H:%M")} GMT', fontweight='bold')
    ax1.set_ylabel('Price')
    ax1.grid(True)
    ax1.legend(loc='upper left', bbox_to_anchor=(1, 1))

    plt.tight_layout()
    plt.subplots_adjust(hspace=0)

    image_filename = f"{symbol}_{event_time.strftime('%Y%m%d_%H%M%S')}.png"
    plt.savefig(image_filename)
    plt.close(fig)

    return image_filename, volatility_unit, max_movement, min_movement, pre_atr_volatility

def process_event(event, date_time_inputs, event_data_list, event_symbols):
    event_summary = {
        'volatility_units': [],
        'max_movements': [],
        'min_movements': [],
        'pre_atr_volatilities': []
    }
    
    event_dates = [
        date for date in sorted(date_time_inputs.keys(), key=lambda x: datetime.strptime(x, '%Y-%m-%d'))
        if any(e['name'] == event['name'] for e in event_data_list[date])
    ]
    
    for idx, date_str in enumerate(event_dates):
        print(f"Processing date: {date_str}")
        
        event_data = next((e for e in event_data_list[date_str] if e['name'] == event['name']), None)
        if not event_data:
            continue
        
        if event['name'] == "Crude Oil Inventories":
            mt5_symbol = event_symbols[event['name']]['mt5']
            polygon_symbol = None
            event_tz = pytz.timezone('Etc/GMT-3')
        else:
            polygon_symbol = event_symbols[event['name']]['polygon']
            mt5_symbol = event_symbols[event['name']]['mt5']
            event_tz = pytz.timezone('Etc/GMT')
        
        event_time_str = event_data["time"]
        event_date = datetime.strptime(date_str, '%Y-%m-%d').date()
        event_time_obj = datetime.strptime(event_time_str, '%H:%M').time()
        
        # Create timezone-aware datetime
        event_time = event_tz.localize(datetime.combine(event_date, event_time_obj))

        start_time = event_time - timedelta(hours=2)
        end_time = event_time + timedelta(hours=2)

        current_spread = get_current_spread(mt5_symbol)
        if current_spread is None:
            current_spread = estimate_typical_spread(mt5_symbol.split('.')[0])

        if event['name'] == "Crude Oil Inventories":
            df = fetch_data(mt5_symbol, start_time, end_time)
        else:
            df = fetch_data(polygon_symbol, start_time, end_time)

        if df is not None and not df.empty:
            is_last_chart = (idx == len(event_dates) - 1)

            event_summary_for_plot = {
                "Avg Volatility Unit": sum(event_summary['volatility_units']) / len(event_summary['volatility_units']) if event_summary['volatility_units'] else 0,
                "Avg Max Movement": sum(event_summary['max_movements']) / len(event_summary['max_movements']) if event_summary['max_movements'] else 0,
                "Avg Min Movement": sum(event_summary['min_movements']) / len(event_summary['min_movements']) if event_summary['min_movements'] else 0,
                "Avg Pre-ATR Volatility": sum(event_summary['pre_atr_volatilities']) / len(event_summary['pre_atr_volatilities']) if event_summary['pre_atr_volatilities'] else 0
            } if is_last_chart else None

            # Convert event_time to match df timezone (UTC)
            event_time_utc = event_time.astimezone(pytz.UTC)
            
            image_filename, volatility_unit, max_movement, min_movement, pre_atr_volatility = plot_ohlc(
                df, polygon_symbol or mt5_symbol, event_time_utc, [event_data],
                atr_multiple, current_spread, atr_period,
                is_last_chart=is_last_chart,
                event_summary=event_summary_for_plot,
                start_time=start_time.astimezone(pytz.UTC),
                end_time=end_time.astimezone(pytz.UTC)
            )

            event_summary['volatility_units'].append(volatility_unit)
            event_summary['max_movements'].append(max_movement / volatility_unit)
            event_summary['min_movements'].append(min_movement / volatility_unit)
            event_summary['pre_atr_volatilities'].append(pre_atr_volatility)

            generated_image_filenames.append(image_filename)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            slide_width = prs.slide_width
            slide_height = prs.slide_height

            left = top = Inches(0)
            width = slide_width
            height = slide_height
            slide.shapes.add_picture(image_filename, left, top, height=height)

        time.sleep(2)
    
    display_event_summary(event, event_summary)
    return event_summary

def display_event_summary(event, summary):
    print(f"Final Summary for {event['name']}:")
    if summary['volatility_units']:
        print(f"Average Volatility Unit: {sum(summary['volatility_units']) / len(summary['volatility_units']):.4f}")
    else:
        print("No data available for Volatility Unit")
    
    if summary['max_movements']:
        print(f"Average Max Movement: {sum(summary['max_movements']) / len(summary['max_movements']):.2f}x")
    else:
        print("No data available for Max Movement")
    
    if summary['min_movements']:
        print(f"Average Min Movement: {sum(summary['min_movements']) / len(summary['min_movements']):.2f}x")
    else:
        print("No data available for Min Movement")
    
    if summary['pre_atr_volatilities']:
        print(f"Average Pre-ATR Volatility: {sum(summary['pre_atr_volatilities']) / len(summary['pre_atr_volatilities']):.4f}")
    else:
        print("No data available for Pre-ATR Volatility")
    print("\n")

# User inputs
atr_multiple = 1.5
atr_period = 14

print("Please paste your Excel-formatted event data (press Enter twice when finished):")
excel_input = ""
while True:
    line = input()
    if line.strip() == "":
        break
    excel_input += line + "\n"

# Parse the Excel input
try:
    date_time_inputs, event_data_list, events, event_symbols = parse_excel_input(excel_input)
except ValueError as e:
    print(f"Error: {e}")
    exit(1)

# Process events
for event in events:
    event_summary = process_event(event, date_time_inputs, event_data_list, event_symbols)

output_pptx_path = "/Users/pawan/Desktop/Quantwater Tech Investments/Research & Development/Research/Blog/Slides/updated_presentation3.pptx"
prs.save(output_pptx_path)

# Clean up temporary image files
for image_filename in generated_image_filenames:
    os.remove(image_filename)

# Shutdown MT5
mt5.shutdown()