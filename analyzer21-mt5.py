import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from datetime import datetime, timedelta
import os
import pytz
import MetaTrader5 as mt5
from mplfinance.original_flavor import candlestick_ohlc
from mpl_toolkits.axes_grid1 import make_axes_locatable
import re
from pptx import Presentation
from pptx.util import Inches

# Initialize MT5 connection
if not mt5.initialize():
    print("MT5 initialization failed")
    mt5.shutdown()

def estimate_typical_spread(symbol):
    major_fx = {
        "EURUSD": 0.00010, "USDJPY": 0.00010, "GBPUSD": 0.00015,
        "USDCHF": 0.00015, "AUDUSD": 0.00015, "USDCAD": 0.00015
    }
    minor_fx = {
        "EURGBP": 0.00020, "EURJPY": 0.00020, "GBPJPY": 0.00025,
        "EURCHF": 0.00025, "AUDJPY": 0.00020, "NZDUSD": 0.00020
    }
    
    if symbol in major_fx:
        return major_fx[symbol]
    elif symbol in minor_fx:
        return minor_fx[symbol]
    else:
        return 0.00030  # Default to 3 pips

def get_current_spread(symbol):
    symbol_info = mt5.symbol_info(symbol)
    if symbol_info is None:
        print(f"Failed to get symbol info for {symbol}")
        return None
    
    spread = symbol_info.ask - symbol_info.bid
    if spread == 0:
        typical_spread = estimate_typical_spread(symbol)
        print(f"MT5 returned 0 spread for {symbol}. Using estimated typical spread: {typical_spread}")
        return typical_spread
    
    return spread

def fetch_data_with_fallback(symbol, start_time, end_time, min_candles=30):
    """
    Attempts to fetch data using progressively larger timeframes if smaller ones fail.
    Ensures a minimum number of candles are available for meaningful analysis.
    """
    timeframes = [
        (mt5.TIMEFRAME_M1, "1 minute"),
        (mt5.TIMEFRAME_M2, "2 minutes"),
        (mt5.TIMEFRAME_M3, "3 minutes"),
        (mt5.TIMEFRAME_M5, "5 minutes"),
        (mt5.TIMEFRAME_M10, "10 minutes"),
        (mt5.TIMEFRAME_M15, "15 minutes"),
        (mt5.TIMEFRAME_M30, "30 minutes"),
        (mt5.TIMEFRAME_H1, "1 hour"),
        (mt5.TIMEFRAME_H4, "4 hours"),
        (mt5.TIMEFRAME_D1, "1 day")
    ]
    
    # Convert to UTC for MT5
    start_time_utc = start_time.astimezone(pytz.UTC)
    end_time_utc = end_time.astimezone(pytz.UTC)
    
    for timeframe, timeframe_name in timeframes:
        print(f"Attempting to fetch {timeframe_name} data for {symbol}...")
        
        # Try to get data from an extended time range
        extended_start = start_time_utc - timedelta(hours=12)  # Extend backward
        extended_end = end_time_utc + timedelta(hours=12)      # Extend forward
        
        rates = mt5.copy_rates_range(symbol, timeframe, extended_start, extended_end)
        
        if rates is not None and len(rates) >= min_candles:
            print(f"Successfully fetched {len(rates)} {timeframe_name} candles")
            df = pd.DataFrame(rates)
            df['time'] = pd.to_datetime(df['time'], unit='s', utc=True)
            
            # Add volume column if it doesn't exist
            if 'volume' not in df.columns:
                df['volume'] = 0
            
            # Trim the data back to the original time range while ensuring minimum candles
            df = df[(df['time'] >= start_time_utc) & (df['time'] <= end_time_utc)]
            
            if len(df) < min_candles:
                print(f"After trimming to event window, only {len(df)} candles remain. Trying larger timeframe...")
                continue
                
            df = calculate_atr(df)
            df.attrs['timeframe'] = timeframe_name
            return df
        else:
            print(f"Insufficient data: Got {len(rates) if rates is not None else 0} candles, need at least {min_candles}")
            
    # If we get here, try one final attempt with the largest available historical data
    print(f"Attempting to fetch maximum available historical data for {symbol}...")
    
    # Try to get a much larger historical range
    extended_start = start_time_utc - timedelta(days=30)  # Go back 30 days
    extended_end = end_time_utc + timedelta(days=30)      # Go forward 30 days
    
    for timeframe, timeframe_name in reversed(timeframes):  # Try from largest timeframe first
        rates = mt5.copy_rates_range(symbol, timeframe, extended_start, extended_end)
        if rates is not None and len(rates) >= min_candles:
            print(f"Successfully fetched {len(rates)} {timeframe_name} candles from extended historical range")
            df = pd.DataFrame(rates)
            df['time'] = pd.to_datetime(df['time'], unit='s', utc=True)
            
            if 'volume' not in df.columns:
                df['volume'] = 0
                
            df = calculate_atr(df)
            df.attrs['timeframe'] = timeframe_name
            return df
    
    print(f"Failed to fetch sufficient data for {symbol} across all timeframes")
    return None

def calculate_atr(df, period=14):
    df['tr1'] = df['high'] - df['low']
    df['tr2'] = (df['high'] - df['close'].shift()).abs()
    df['tr3'] = (df['low'] - df['close'].shift()).abs()
    df['tr'] = df[['tr1', 'tr2', 'tr3']].max(axis=1)
    df['atr'] = df['tr'].rolling(window=period).mean()
    return df

def parse_excel_input(excel_input):
    excel_input = excel_input.replace('^I', '\t')
    lines = excel_input.strip().split('\n')
    events = []
    date_time_inputs = {}
    event_data_list = {}
    event_symbols = {}

    current_event_name = ""
    current_event_time = ""

    for line in lines:
        if line.startswith('Event'):
            parts = line.split('\t')
            current_event_name = parts[1].strip()
            events.append({'name': current_event_name})
        elif line.startswith('Actual:'):
            parts = line.split('\t')
            symbol = parts[-1].strip()
            event_symbols[current_event_name] = symbol
        elif line.startswith('Time (GMT):'):
            parts = line.split('\t')
            current_event_time = parts[1].strip()
        elif re.match(r'\d{2}-\w{3}-\d{2}', line):
            parts = line.split('\t')
            date_str = parts[0].strip()
            date = datetime.strptime(date_str, '%d-%b-%y').strftime('%Y-%m-%d')
            actual = parts[1].strip().replace('%', '')
            forecast = parts[2].strip().replace('%', '') or "NA"
            previous = parts[3].strip().replace('%', '') or "NA"
            time_str = parts[4].strip()

            if date not in date_time_inputs:
                date_time_inputs[date] = {
                    "date": date,
                    "event": current_event_time + "3:00",
                }
                event_data_list[date] = []

            event_data_list[date].append({
                "name": current_event_name,
                "actual": actual,
                "forecast": forecast,
                "previous": previous,
                "time": time_str
            })

    return date_time_inputs, event_data_list, events, event_symbols

# Initialize presentation
template_path = "/Users/pawan/Desktop/Quantwater Tech Investments/Research & Development/Research/Blog/Slides/template3.pptx"
prs = Presentation(template_path)
generated_image_filenames = []

def plot_ohlc(df, symbol, event_time, event_data, atr_multiple, spread, atr_period=14,
              is_last_chart=False, event_summary=None, start_time=None, end_time=None):
    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(14.5, 9),
                                   gridspec_kw={'height_ratios': [4, 1]}, sharex=True)
    
    df['time_num'] = mdates.date2num(df['time'])
    ohlc = df[['time_num', 'open', 'high', 'low', 'close']].values

    plt.style.use('ggplot')

    # Plot candlesticks
    candlestick_ohlc(ax1, ohlc, width=0.0004, colorup='green', colordown='red')

    # Configure x-axis with proper timeline
    ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M'))
    ax1.xaxis.set_major_locator(mdates.HourLocator(interval=1))
    ax1.xaxis.set_minor_locator(mdates.MinuteLocator(byminute=[15, 30, 45]))
    plt.setp(ax1.get_xticklabels(), rotation=45, ha='right')

    # Add grid with minor lines
    ax1.grid(True, which='major', linestyle='-', alpha=0.6)
    ax1.grid(True, which='minor', linestyle='--', alpha=0.3)

    closest_index = (df['time'] - event_time).abs().idxmin()
    closest_time = df.loc[closest_index, 'time']
    event_price = df.loc[closest_index, 'open']

    ax1.plot(closest_time, event_price, 'X', color='red', markersize=10, label='Event Time')

    atr_value = df.loc[closest_index, 'atr']
    volatility_unit = atr_value * atr_multiple + spread

    # Plot volatility units
    for i in range(1, 6):
        upper_range = event_price + i * volatility_unit
        lower_range = event_price - i * volatility_unit
        ax1.axhline(upper_range, color='blue', linestyle='--', alpha=0.5, label=f'+{i} Volatility Unit' if i == 1 else "")
        ax1.axhline(lower_range, color='orange', linestyle='--', alpha=0.5, label=f'-{i} Volatility Unit' if i == 1 else "")

    post_event_df = df[df['time'] >= event_time]
    max_price = post_event_df['high'].max()
    min_price = post_event_df['low'].min()
    max_movement = max_price - event_price
    min_movement = event_price - min_price

    if not post_event_df[post_event_df['high'] == max_price].empty:
        max_time = post_event_df.loc[post_event_df['high'] == max_price, 'time'].iloc[0]
        ax1.plot(max_time, max_price, 'gx', markersize=10, label='Max Peak')
    
    if not post_event_df[post_event_df['low'] == min_price].empty:
        min_time = post_event_df.loc[post_event_df['low'] == min_price, 'time'].iloc[0]
        ax1.plot(min_time, min_price, 'rx', markersize=10, label='Min Peak')

    pre_atr_df = df[df['time'] < event_time].iloc[-(atr_period*2):-atr_period]
    pre_atr_volatility = pre_atr_df['high'].max() - pre_atr_df['low'].min()

    # Plot volume with color coding
    volume_colors = ['green' if close >= open else 'red' 
                    for open, close in zip(df['open'], df['close'])]
    ax2.bar(df['time'], df['volume'], width=0.0004, color=volume_colors, alpha=0.7)
    ax2.set_ylabel('Volume')
    
    # Add grid to volume subplot
    ax2.grid(True, alpha=0.3)

    # Include timeframe in the information text
    info_text = (
        f'Timeframe: {df.attrs["timeframe"]}\n'
        f'ATR: {atr_value:.4f}\n'
        f'Spread: {spread:.5f}\n'
        f'Pre-ATR Volatility: {pre_atr_volatility:.4f}\n'
        f'\n'
        f'Volatility Unit (ATR*{atr_multiple}+spread): {volatility_unit:.4f}\n'
        f'Max Movement: {max_movement/volatility_unit:.2f}x\n'
        f'Min Movement: {min_movement/volatility_unit:.2f}x'
    )

    if is_last_chart and event_summary and event_summary['Avg Volatility Unit'] > 0:
        info_text += (
            f'\n\nEvent Summary:\n'
            f'Avg Volatility Unit: {event_summary["Avg Volatility Unit"]:.4f}\n'
            f'Avg Max Movement: {event_summary["Avg Max Movement"]:.2f}x\n'
            f'Avg Min Movement: {event_summary["Avg Min Movement"]:.2f}x\n'
            f'Avg Pre-ATR Volatility: {event_summary["Avg Pre-ATR Volatility"]:.4f}'
        )

    ax1.text(1.01, 0.5, info_text, transform=ax1.transAxes, ha='left', va='center', fontsize=9)

    if start_time and end_time:
        ax1.set_xlim(mdates.date2num(start_time), mdates.date2num(end_time))

    divider = make_axes_locatable(ax2)
    ax_event = divider.append_axes("bottom", size="15%", pad=0.1, sharex=ax1)
    ax_event.axis('off')

    plt.subplots_adjust(bottom=0.45)  

    for i, event in enumerate(event_data):
        event_name = event['name']
        actual = event['actual']
        forecast = event['forecast']
        previous = event['previous']

        if forecast == "NA":
            event_color = 'black'
        else:
            try:
                if float(actual) > float(forecast):
                    event_color = 'green'
                elif float(actual) < float(forecast):
                    event_color = 'red'
                else:
                    event_color = 'blue'
            except ValueError:
                event_color = 'black'

        event_time_num = mdates.date2num(event_time)
        event_text_line = f"{event_name} | Actual: {actual} | Forecast: {forecast} | Previous: {previous}"
        y_offset = -0.2 - i * 0.1  
        ax_event.text(event_time_num, y_offset, event_text_line, color=event_color,
                      ha='center', va='top', fontsize=9, rotation=0, transform=ax_event.transData, fontweight='bold')

    ax1.set_title(f'{symbol} Price Action - {event_time.strftime("%Y-%m-%d %H:%M")} GMT', fontweight='bold')
    ax1.set_ylabel('Price')
    ax1.legend(loc='upper left', bbox_to_anchor=(1, 1))

    plt.tight_layout()
    plt.subplots_adjust(hspace=0)

    image_filename = f"{symbol}_{event_time.strftime('%Y%m%d_%H%M%S')}.png"
    plt.savefig(image_filename, bbox_inches='tight')
    plt.close(fig)

    return image_filename, volatility_unit, max_movement, min_movement, pre_atr_volatility

def process_event(event, date_time_inputs, event_data_list, event_symbols):
    event_summary = {
        'volatility_units': [],
        'max_movements': [],
        'min_movements': [],
        'pre_atr_volatilities': []
    }
    
    event_dates = [
        date for date in sorted(date_time_inputs.keys(), key=lambda x: datetime.strptime(x, '%Y-%m-%d'))
        if any(e['name'] == event['name'] for e in event_data_list[date])
    ]
    
    for idx, date_str in enumerate(event_dates):
        print(f"\nProcessing date: {date_str}")
        
        event_data = next((e for e in event_data_list[date_str] if e['name'] == event['name']), None)
        if not event_data:
            continue
            
        try:
            mt5_symbol = event_symbols[event['name']]
            event_tz = pytz.timezone('Etc/GMT')
            
            event_time_str = event_data["time"]
            event_date = datetime.strptime(date_str, '%Y-%m-%d').date()
            event_time_obj = datetime.strptime(event_time_str, '%H:%M').time()
            
            # Create timezone-aware datetime
            event_time = event_tz.localize(datetime.combine(event_date, event_time_obj))

            # Extended time range
            start_time = event_time - timedelta(hours=8)  # Extended further
            end_time = event_time + timedelta(hours=8)    # Extended further

            current_spread = get_current_spread(mt5_symbol)
            if current_spread is None:
                current_spread = estimate_typical_spread(mt5_symbol)

            df = fetch_data_with_fallback(mt5_symbol, start_time, end_time)

            if df is not None and not df.empty and len(df) >= 30:  # Ensure minimum data points
                print(f"Using {df.attrs['timeframe']} timeframe data with {len(df)} candles")
                is_last_chart = (idx == len(event_dates) - 1)

                event_summary_for_plot = {
                    "Avg Volatility Unit": sum(event_summary['volatility_units']) / len(event_summary['volatility_units']) if event_summary['volatility_units'] else 0,
                    "Avg Max Movement": sum(event_summary['max_movements']) / len(event_summary['max_movements']) if event_summary['max_movements'] else 0,
                    "Avg Min Movement": sum(event_summary['min_movements']) / len(event_summary['min_movements']) if event_summary['min_movements'] else 0,
                    "Avg Pre-ATR Volatility": sum(event_summary['pre_atr_volatilities']) / len(event_summary['pre_atr_volatilities']) if event_summary['pre_atr_volatilities'] else 0
                } if is_last_chart else None

                event_time_utc = event_time.astimezone(pytz.UTC)
                
                image_filename, volatility_unit, max_movement, min_movement, pre_atr_volatility = plot_ohlc(
                    df, mt5_symbol, event_time_utc, [event_data],
                    atr_multiple, current_spread, atr_period,
                    is_last_chart=is_last_chart,
                    event_summary=event_summary_for_plot,
                    start_time=start_time.astimezone(pytz.UTC),
                    end_time=end_time.astimezone(pytz.UTC)
                )

                if all(x is not None for x in [volatility_unit, max_movement, min_movement, pre_atr_volatility]):
                    event_summary['volatility_units'].append(volatility_unit)
                    event_summary['max_movements'].append(max_movement / volatility_unit)
                    event_summary['min_movements'].append(min_movement / volatility_unit)
                    event_summary['pre_atr_volatilities'].append(pre_atr_volatility)

                generated_image_filenames.append(image_filename)

                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)

                slide_width = prs.slide_width
                slide_height = prs.slide_height

                left = top = Inches(0)
                width = slide_width
                height = slide_height
                slide.shapes.add_picture(image_filename, left, top, height=height)
            else:
                print(f"Insufficient data for {mt5_symbol} on {date_str}")
        except Exception as e:
            print(f"Error processing event for {date_str}: {str(e)}")
            continue

    display_event_summary(event, event_summary)
    return event_summary

def display_event_summary(event, summary):
    print(f"\nFinal Summary for {event['name']}:")
    if summary['volatility_units']:
        print(f"Average Volatility Unit: {sum(summary['volatility_units']) / len(summary['volatility_units']):.4f}")
    else:
        print("No data available for Volatility Unit")
    
    if summary['max_movements']:
        print(f"Average Max Movement: {sum(summary['max_movements']) / len(summary['max_movements']):.2f}x")
    else:
        print("No data available for Max Movement")
    
    if summary['min_movements']:
        print(f"Average Min Movement: {sum(summary['min_movements']) / len(summary['min_movements']):.2f}x")
    else:
        print("No data available for Min Movement")
    
    if summary['pre_atr_volatilities']:
        print(f"Average Pre-ATR Volatility: {sum(summary['pre_atr_volatilities']) / len(summary['pre_atr_volatilities']):.4f}")
    else:
        print("No data available for Pre-ATR Volatility")
    print("\n")

# User inputs
atr_multiple = 1.5
atr_period = 14

print("Please paste your Excel-formatted event data (press Enter twice when finished):")
print("Format example:")
print("Event\tNFP")
print("Actual:\tEURUSD")
print("Time (GMT):\t13:30")
print("01-Mar-24\t275K\t198K\t229K\t13:30")
print("\n")

excel_input = ""
while True:
    line = input()
    if line.strip() == "":
        break
    excel_input += line + "\n"

# Parse the Excel input
try:
    date_time_inputs, event_data_list, events, event_symbols = parse_excel_input(excel_input)
except ValueError as e:
    print(f"Error: {e}")
    exit(1)

# Process events
for event in events:
    event_summary = process_event(event, date_time_inputs, event_data_list, event_symbols)

output_pptx_path = "/Users/pawan/Desktop/Quantwater Tech Investments/Research & Development/Research/Blog/Slides/updated_presentation3.pptx"
prs.save(output_pptx_path)

# Clean up temporary image files
for image_filename in generated_image_filenames:
    try:
        os.remove(image_filename)
    except Exception as e:
        print(f"Error removing temporary file {image_filename}: {e}")

# Shutdown MT5
mt5.shutdown()